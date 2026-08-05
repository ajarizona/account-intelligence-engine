"""
render_deck.py

Turns an engine result (from engine.score_company) into a meeting-ready .pptx.

Design intent:
  - The deck is for the room, not for the analyst. Rubric scores are decision
    support, so the narrative slides carry no numbers. The scores live on one
    appendix slide, for the "why these" question.
  - Structure follows the crawl, walk, run spine the engine already produces.
  - The honest slides (what execution takes, how we measure) stay in, because
    they are what separate a readout from a vendor pitch.

Dependency: python-pptx only, to keep the engine a single-language project.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy
from lxml import etree

# ----------------------------------------------------------------------------
# Palette (Midnight Executive, with a gold accent) and type
# ----------------------------------------------------------------------------
NAVY = RGBColor(0x1E, 0x27, 0x61)
NAVY_DEEP = RGBColor(0x16, 0x1C, 0x45)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
GOLD = RGBColor(0xE8, 0xB0, 0x4B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x33, 0x38, 0x50)      # body text on light
MUTED = RGBColor(0x6B, 0x71, 0x8C)    # captions
CARD = RGBColor(0xF2, 0xF5, 0xFC)     # light card tint
CARD_LINE = RGBColor(0xDD, 0xE4, 0xF5)

HEAD_FONT = "Cambria"
BODY_FONT = "Calibri"

EMU = 914400
SW, SH = 13.333, 7.5

PHASE_INTRO = {
    "crawl": "Prove governed value on a contained use case, where the risk is lowest and the audit story is cleanest.",
    "walk": "Extend into service and personalization once the foundation holds.",
    "run": "Turn on the growth engine, on a foundation the risk side already trusts.",
}


# ----------------------------------------------------------------------------
# Low-level helpers
# ----------------------------------------------------------------------------

def _slide(prs, bg=WHITE):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SW), Inches(SH))
    r.fill.solid()
    r.fill.fore_color.rgb = bg
    r.line.fill.background()
    _no_shadow(r)
    # send background to back
    sp = r._element
    sp.getparent().remove(sp)
    s.shapes._spTree.insert(2, sp)
    return s


def _no_shadow(shape):
    spPr = shape._element.spPr
    for tag in ("a:effectLst",):
        for e in spPr.findall(f'{{http://schemas.openxmlformats.org/drawingml/2006/main}}effectLst'):
            spPr.remove(e)
    el = etree.SubElement(spPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst')


def _soft_shadow(shape):
    """Add a subtle drop shadow (not an edge stripe)."""
    spPr = shape._element.spPr
    ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    for e in spPr.findall(f'{ns}effectLst'):
        spPr.remove(e)
    eff = etree.SubElement(spPr, f'{ns}effectLst')
    shdw = etree.SubElement(eff, f'{ns}outerShdw')
    shdw.set('blurRad', '90000')
    shdw.set('dist', '38100')
    shdw.set('dir', '5400000')
    shdw.set('rotWithShape', '0')
    clr = etree.SubElement(shdw, f'{ns}srgbClr')
    clr.set('val', '1E2761')
    alpha = etree.SubElement(clr, f'{ns}alpha')
    alpha.set('val', '18000')


def _text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          space_after=6):
    """runs: list of paragraphs, each a list of (text, font, size, color, bold, italic)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for (txt, font, size, color, bold, italic) in para:
            r = p.add_run()
            r.text = txt
            r.font.name = font
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.italic = italic
    return tb


def _card(slide, x, y, w, h, tint=CARD, line=CARD_LINE, shadow=True):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.adjustments[0] = 0.06
    c.fill.solid()
    c.fill.fore_color.rgb = tint
    c.line.color.rgb = line
    c.line.width = Pt(1)
    if shadow:
        _soft_shadow(c)
    else:
        _no_shadow(c)
    return c


def _pill(slide, x, y, w, h, text, fill, txt_color):
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    p.adjustments[0] = 0.5
    p.fill.solid()
    p.fill.fore_color.rgb = fill
    p.line.fill.background()
    _no_shadow(p)
    tf = p.text_frame
    tf.margin_top = 0
    tf.margin_bottom = 0
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = text
    run.font.name = BODY_FONT
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = txt_color
    return p


def _num_circle(slide, x, y, d, n, fill=NAVY, txt=WHITE):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    c.line.fill.background()
    _no_shadow(c)
    tf = c.text_frame
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(n)
    r.font.name = HEAD_FONT
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = txt
    return c


def _title(slide, text, color=NAVY, y=0.55):
    _text(slide, 0.7, y, SW - 1.4, 0.9,
          [[(text, HEAD_FONT, 34, color, True, False)]])


def _kicker(slide, text, color=GOLD, y=0.42):
    _text(slide, 0.72, y, SW - 1.4, 0.3,
          [[(text.upper(), BODY_FONT, 12, color, True, False)]])


# ----------------------------------------------------------------------------
# Slides
# ----------------------------------------------------------------------------

def _slide_title(prs, result):
    s = _slide(prs, bg=NAVY)
    company = result["company"]["name"]
    vp = result.get("value_pattern")
    thesis = vp["thesis"] if vp else "A grounded, sequenced plan built from the company's own priorities."
    _text(s, 0.9, 1.5, SW - 1.8, 0.4, [[("ACCOUNT INTELLIGENCE ENGINE", BODY_FONT, 13, GOLD, True, False)]])
    _text(s, 0.9, 2.0, SW - 1.8, 1.4,
          [[(company, HEAD_FONT, 46, WHITE, True, False)]])
    _text(s, 0.9, 3.4, SW - 1.8, 0.6,
          [[("Agentic Value Readout", HEAD_FONT, 24, ICE, False, True)]])
    _text(s, 0.9, 4.4, SW - 3.2, 1.4,
          [[(thesis, BODY_FONT, 15, ICE, False, False)]])
    note = result["company"].get("source_note")
    if note:
        _text(s, 0.9, 6.6, SW - 1.8, 0.5,
              [[(note, BODY_FONT, 10, RGBColor(0x9A, 0xA8, 0xD0), False, True)]])
    return s


def _slide_what_we_heard(prs, result):
    s = _slide(prs)
    _kicker(s, "What we heard")
    _title(s, "Your stated priorities, in your own words")
    stated = sorted(result.get("stated_objectives", []),
                    key=lambda o: o.get("emphasis", 0), reverse=True)[:6]
    # two columns of priority cards
    cols, gap = 2, 0.4
    cw = (SW - 1.4 - gap) / cols
    ch = 1.15
    x0, y0 = 0.7, 1.7
    for i, o in enumerate(stated):
        col = i % cols
        row = i // cols
        x = x0 + col * (cw + gap)
        y = y0 + row * (ch + 0.25)
        _card(s, x, y, cw, ch)
        _text(s, x + 0.25, y + 0.18, cw - 0.5, 0.4,
              [[(o["objective_name"], BODY_FONT, 15, NAVY, True, False)]])
        ev = o.get("evidence", "")
        _text(s, x + 0.25, y + 0.55, cw - 0.5, 0.5,
              [[(ev, BODY_FONT, 11, INK, False, False)]], space_after=0)
    return s


def _slide_opportunity(prs, result):
    s = _slide(prs)
    _kicker(s, "The opportunity")
    _title(s, "Where we can move first")
    top = result["ranked_use_cases"][:4]
    y = 1.75
    rh = 1.15
    for i, u in enumerate(top, 1):
        _num_circle(s, 0.75, y + 0.12, 0.55, i, fill=NAVY)
        head = f"{u['solution_name']} for {u['primary_objective_name']}"
        _text(s, 1.55, y, SW - 2.4, 0.4,
              [[(head, BODY_FONT, 16, NAVY, True, False)]])
        body = u["rationale"]
        if u.get("also_supports"):
            body = body  # rationale already includes also-supports
        _text(s, 1.55, y + 0.42, SW - 2.6, 0.6,
              [[(body, BODY_FONT, 12, INK, False, False)]], space_after=0)
        y += rh
    return s


def _slide_roadmap(prs, result):
    s = _slide(prs)
    _kicker(s, "The roadmap")
    _title(s, "Crawl, walk, run")
    phases = ("crawl", "walk", "run")
    gap = 0.4
    cw = (SW - 1.4 - 2 * gap) / 3
    x0, y0 = 0.7, 1.75
    ch = 3.7
    fills = {"crawl": NAVY, "walk": RGBColor(0x3B, 0x47, 0x8C), "run": GOLD}
    txt_on = {"crawl": WHITE, "walk": WHITE, "run": NAVY_DEEP}
    for i, ph in enumerate(phases):
        x = x0 + i * (cw + gap)
        _card(s, x, y0, cw, ch)
        _pill(s, x + 0.25, y0 + 0.25, 1.5, 0.42, ph.capitalize(), fills[ph], txt_on[ph])
        _text(s, x + 0.25, y0 + 0.85, cw - 0.5, 0.9,
              [[(PHASE_INTRO.get(ph, ""), BODY_FONT, 11.5, INK, False, False)]])
        items = result["roadmap"].get(ph, [])
        yy = y0 + 1.85
        for u in items:
            _text(s, x + 0.25, yy, cw - 0.5, 0.35,
                  [[(u["solution_name"], BODY_FONT, 13, NAVY, True, False)]], space_after=0)
            _text(s, x + 0.25, yy + 0.28, cw - 0.5, 0.35,
                  [[(u["primary_objective_name"], BODY_FONT, 10.5, MUTED, False, True)]], space_after=0)
            yy += 0.72
    # tensions strip
    tensions = result.get("strategic_tensions", [])[:3]
    if tensions:
        names = "   |   ".join(t["name"] for t in tensions)
        _text(s, 0.7, y0 + ch + 0.25, SW - 1.4, 0.6,
              [[("We sequence to respect these tensions:  ", BODY_FONT, 11, NAVY, True, False),
                (names, BODY_FONT, 11, MUTED, False, True)]])
    return s


def _slide_challenges(prs, result):
    s = _slide(prs)
    _kicker(s, "What execution takes")
    _title(s, "An honest look at delivery")
    challenges = result.get("execution_challenges", [])[:5]
    cols, gap = 2, 0.4
    cw = (SW - 1.4 - gap) / cols
    ch = 1.35
    x0, y0 = 0.7, 1.7
    for i, cobj in enumerate(challenges):
        col = i % cols
        row = i // cols
        x = x0 + col * (cw + gap)
        y = y0 + row * (ch + 0.22)
        _card(s, x, y, cw, ch)
        _text(s, x + 0.25, y + 0.16, cw - 0.5, 0.4,
              [[(cobj["name"], BODY_FONT, 13.5, NAVY, True, False)]])
        _text(s, x + 0.25, y + 0.6, cw - 0.5, 0.7,
              [[(cobj["description"], BODY_FONT, 10.5, INK, False, False)]], space_after=0)
    return s


def _slide_measure(prs, result):
    s = _slide(prs, bg=NAVY)
    _text(s, 0.9, 0.9, SW - 1.8, 0.4, [[("HOW WE MEASURE", BODY_FONT, 13, GOLD, True, False)]])
    _text(s, 0.9, 1.4, SW - 1.8, 0.9, [[("We prove value, we do not assert it", HEAD_FONT, 32, WHITE, True, False)]])
    points = [
        ("Baseline first", "We record the starting point before anything ships, so the gain is measured against reality."),
        ("Claim lift against a holdout", "We report incremental lift versus a control, not gross results."),
        ("Agree the scorecard with Finance", "We set the measures up front, so the review confirms value rather than negotiating it."),
    ]
    y = 2.7
    for i, (h, b) in enumerate(points, 1):
        _num_circle(s, 0.9, y + 0.05, 0.55, i, fill=GOLD, txt=NAVY_DEEP)
        _text(s, 1.7, y, SW - 3.0, 0.4, [[(h, BODY_FONT, 17, ICE, True, False)]])
        _text(s, 1.7, y + 0.42, SW - 3.2, 0.6, [[(b, BODY_FONT, 12.5, RGBColor(0xC7, 0xD1, 0xEE), False, False)]], space_after=0)
        y += 1.15
    return s


def _slide_appendix(prs, result):
    s = _slide(prs)
    _kicker(s, "Appendix")
    _title(s, "Ranked use cases and scores")
    _text(s, 0.72, 1.4, SW - 1.4, 0.4,
          [[("Decision support, not a customer metric. Rubric: alignment 0.35, value 0.30, feasibility 0.20, evidence 0.15.",
             BODY_FONT, 11, MUTED, False, True)]])
    rows = result["ranked_use_cases"]
    n = len(rows) + 1
    tbl_shape = s.shapes.add_table(n, 6, Inches(0.7), Inches(1.9), Inches(SW - 1.4), Inches(0.5 * n))
    table = tbl_shape.table
    headers = ["Solution", "Primary objective", "Align", "Value", "Feasible", "Score"]
    widths = [3.2, 3.5, 1.3, 1.3, 1.3, 1.3]
    for j, w in enumerate(widths):
        table.columns[j].width = Inches(w)
    for j, htext in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        r = p.add_run(); r.text = htext
        r.font.name = BODY_FONT; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT if j < 2 else PP_ALIGN.CENTER
    for i, u in enumerate(rows, 1):
        sc = u["scores"]
        vals = [u["solution_name"], u["primary_objective_name"],
                f"{sc['strategic_alignment']:.2f}", f"{sc['value_at_stake']:.2f}",
                f"{sc['time_to_value']:.2f}", f"{u['weighted_total']:.3f}"]
        for j, v in enumerate(vals):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 else CARD
            p = cell.text_frame.paragraphs[0]
            r = p.add_run(); r.text = v
            r.font.name = BODY_FONT; r.font.size = Pt(10.5)
            r.font.color.rgb = NAVY if j == 5 else INK
            r.font.bold = (j == 5)
            p.alignment = PP_ALIGN.LEFT if j < 2 else PP_ALIGN.CENTER
    return s


def _slide_close(prs, result):
    s = _slide(prs, bg=NAVY_DEEP)
    _text(s, 0.9, 2.7, SW - 1.8, 0.9, [[("The next step", HEAD_FONT, 34, WHITE, True, False)]])
    _text(s, 0.9, 3.7, SW - 3.0, 1.2,
          [[("Prove the crawl phase on one contained use case, earn the trust, then let the roadmap carry the growth.",
             BODY_FONT, 16, ICE, False, False)]])
    return s


# ----------------------------------------------------------------------------
# Entry point
# ----------------------------------------------------------------------------

def build_deck(result, out_path):
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)

    _slide_title(prs, result)
    if result.get("stated_objectives"):
        _slide_what_we_heard(prs, result)
    _slide_opportunity(prs, result)
    _slide_roadmap(prs, result)
    if result.get("execution_challenges"):
        _slide_challenges(prs, result)
    _slide_measure(prs, result)
    _slide_appendix(prs, result)
    _slide_close(prs, result)

    prs.save(out_path)
    return out_path
